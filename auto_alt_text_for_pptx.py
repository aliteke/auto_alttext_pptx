"""
    Extracts images from a PPTX file, writes them into the same directory.
    Request GoogleVertexAI Image Captioning Service to generate a image caption for each image extracted from PPTX file in the directory.
    Author: ali.tekeoglu@jhu.edu
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from tqdm import tqdm
import argparse
import os
import sys
import base64
import requests
import json
import time
import csv

# TODO: Take a look at this repository: https://github.com/waltervanheuven/auto-alt-text
# TODO: Tutorial https://cloud.google.com/vertex-ai/generative-ai/docs/model-reference/image-captioning
SLEEP_TIME = 4

def getCaptionsForDir(dirName, bearer, outputFile):
    """ Chek if output file exists. If exists; find the images that doesn't have a caption listed in the outputFile, and get the caption for them."""
    lst_existing_captions=list()
    lst_images = [y for y in os.listdir(dirName) if "png" in y.lower() or "jpg" in y.lower() or "jpeg" in y.lower() or "gif" in y.lower() or "tiff" in y.lower()]

    if os.path.exists(outputFile):
        with open(outputFile, 'r') as file:
            lines = csv.reader(file)
            for l in lines:
                lst_existing_captions.append(os.path.basename(l[0]))

    lst_missing_captions = list(set(lst_images) - set(lst_existing_captions))
    print(f"[+] Missing captions for: {lst_missing_captions}")
    for image in tqdm(lst_missing_captions):
        print(f"[+] Sleeping {SLEEP_TIME} seconds....")
        time.sleep(SLEEP_TIME)
        print(f"[+] Getting caption for {image}")
        getCaption(os.path.join(dirName, image), bearer, outputFile)


def getCaption(imageName, bearer, outputFile): 
    """ we got the authorizaion bearer through Google-Cloud-Console command:
        aliteke@cloudshell:~/images (rising-analogy-272214)$ gcloud auth print-access-token
        [!] imageName should be the full path to the image file!
        [!] outputFile should be the full path to the output file!
    """
    projectNameID = "rising-analogy-272214"
    location = "us-central1"

    headers = {"Authorization": f"Bearer {bearer}", "Content-Type": "application/json; charset=utf-8"}
    
    response_count = 1
    language_code = "en"
    b64_image = getB64forImage(imageName)
    data = { "instances": [ { "image": { "bytesBase64Encoded": b64_image } } ], "parameters": { "sampleCount": response_count, "language": language_code } }

    url = f"https://{location}-aiplatform.googleapis.com/v1/projects/{projectNameID}/locations/{location}/publishers/google/models/imagetext:predict"
   
    response = requests.post(url, json=data, headers=headers)
    print(f"[+] Request: url: {url}, headers: {headers}")
    print(f"[+] Status: {response.status_code}, Caption for {imageName}: {response.json()}")
    if response.status_code == 429:
        print(f"[+] Recursively calling itself again in a {SLEEP_TIME*5} seconds...")
        time.sleep(SLEEP_TIME*5)
        getCaption(imageName, bearer, outputFile)
    elif response.status_code == 401:
        print(f"[!] Response: {response.json()}")
        sys.exit(f"[!] 401: Something wrong with Bearer token...")
        
    elif response.status_code == 200:
        if len(response.json())>=1:
            with open(outputFile, 'a') as f:
                f.write(f"{imageName},\"{response.json()['predictions'][0]}\"\n")
        elif len(response.json())==0:
            with open(outputFile, 'a') as f:
                f.write(f"{imageName},\"not available:\"\n")

    
def getB64forImage(imageName):
    with open(f"{imageName}", "rb") as image_file:
        data = base64.b64encode(image_file.read()).decode('utf-8')
    return data


def writeImageToFile(dirName, pageNum, shapeNum, imageBytes, extension):
    fname = f"{dirName}/image_pg{pageNum}_idx{shapeNum}.{extension}"
    print(f"[+] Writing image to file: {fname}") 
    
    if not os.path.exists(fname):
        with open(fname, "wb") as f:
            f.write(imageBytes)
    else:
        print("[!] File already exists, skipping")


def shape_alt_text(shape: BaseShape) -> str:
    """Alt-text defined in shape's `descr` attribute, or "" if not present."""
    return shape._element._nvXxPr.cNvPr.attrib.get("descr", "")


def shape_alt_text_set(shape: BaseShape, alt_text: str):
    """Change the Alt-text that is defined in shape's `descr` attribute """
    shape._element._nvXxPr.cNvPr.descr = alt_text
    #print(f'[+] AFTER INSERT to DICTIONARY: {shape._element._nvXxPr.cNvPr.items()}')


def listImageAlttexts(pptxFname):
    prs = Presentation(pptxFname.replace("\\", ""))
    for slideNum, slide in enumerate(prs.slides):
        for shapeNum, shape in enumerate([s for s in slide.shapes if hasattr(s, 'image')]):
            print(f"Slide#: {slideNum}, Shape# {shapeNum}, SlideID: {slide.slide_id}, Shape_type: {shape.shape_type}, Alt Text: {shape_alt_text(shape)}, ImageEXT: {shape.image.ext}")


def getAltTextFromPredsFile(genCaptionsFname, imageFname):
    preds=list()
    with open(genCaptionsFname, 'r') as file:
        lines = csv.reader(file)
        for l in lines:
            preds.append(l)

    for item in preds:
        fname =  item[0]
        pred = item[1]
        if imageFname in fname:
            return pred


def updateImageAlttexts(pptxFname, genCaptionsFname):
    print(f"[+] Updating image alt text for ppt file: {pptxFname} with captions from: {genCaptionsFname}")
    prs = Presentation(pptxFname)
    for slideNum, slide in enumerate(prs.slides):
        #for shapeNum, shape in enumerate([s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]):
        for shapeNum, shape in enumerate([s for s in slide.shapes if hasattr(s, 'image')]):
            imageFname = f"image_pg{slideNum}_idx{shapeNum}.{shape.image.ext}"
            alt_text = getAltTextFromPredsFile(genCaptionsFname, imageFname)
            print(f"[+] Setting AltText: {alt_text} for {imageFname}")
            shape_alt_text_set(shape, alt_text)

    prs.save(pptxFname)


def resetImageAlttexts(pptxFname):
    """Resets all the Alt-texts to empty string for the images in the pptx file."""
    prs = Presentation(pptxFname)
    for slideNum, slide in enumerate(prs.slides):
        #for shapeNum, shape in enumerate([s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]):
        for shapeNum, shape in enumerate([s for s in slide.shapes if hasattr(s, 'image')]):
            imageFname = f"image_pg{slideNum}_idx{shapeNum}.{shape.image.ext}"
            print(f"[+] Resetting AltText to empty string for {imageFname}")
            shape_alt_text_set(shape, "")
    
    prs.save(pptxFname)


def extractImagesFromPPTX(fname):
    prs = Presentation(fname)
    for slideNum, slide in enumerate(prs.slides):
        #for shapeNum, shape in enumerate([s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]):
        for shapeNum, shape in enumerate([s for s in slide.shapes if hasattr(s, 'image')]):
            #print(f"Slide#: {slideNum}, Shape# {shapeNum}, SlideID: {slide.slide_id}, Shape_type: {shape.shape_type}")
            image_bytes = shape.image.blob    # image's file contents
            #shape_alt_text_set(shape, "Test Alternative Text ALI...")
            dirName=os.path.dirname(fname)
            writeImageToFile(dirName, slideNum, shapeNum, image_bytes, shape.image.ext)

    # TODO: Save the changes to alt-text
    #prs.save(fname)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract images from PPTX")
    parser.add_argument('-p', '--pptx', required=False, type=str, help=".pptx filename to extract images")
    parser.add_argument('-i', '--image', required=False, type=str, help=".png filename to get a caption from GoogleCloudVertex")
    parser.add_argument('-d', '--dir', required=False, type=str, help="a directory of .png images (extracted from a .pptx file) to get the caption for.")
    parser.add_argument('-b', '--bearer', required=False, help="GoogleCloud VertexAI auth token, can be obtained by running `$ gcloud auth print-access-token` on the cloudshell")
    parser.add_argument('-o', '--output', required=False, type=str, help="When using -dir or -image, we need to provide an output file to write the caption responses from GoogleCloud")
    parser.add_argument('-u', '--update', required=False, type=str, help="to update the Alternative Text for the images. Need to provide a csv file generated with --output in a prev run from --dir or --image")
    parser.add_argument('-l', '--list', action='store_true', required=False, help="list all the images and their Alternative Texts in the pptx file")
    parser.add_argument('-r', '--reset', action='store_true', required=False, help="if turned on, resets all the Alternative Texts to empty string in the pptx file")
    
    args = parser.parse_args()

    if args.pptx and args.image and args.dir:
        sys.exit(f"[!] Only provide one flag; either PNG image, PPTX file, or a directory")

    if args.pptx and args.reset:
        resetImageAlttexts(args.pptx)
        sys.exit(f"[+] Finished, reseting all the Alt-texts to empty string in the pptx file {args.pptx}")

    if args.pptx and not args.update and not args.list:
        if not os.path.exists(f"{args.pptx}"):
            sys.exit(f"[!] {args.pptx} doesn't exist")
        extractImagesFromPPTX(args.pptx)

    elif args.image and args.bearer and args.output:
        if not os.path.exists(f"{args.image}") or not os.path.exists(f"{args.output}"):
            sys.exit(f"[!] {args.image} or {args.output} doesn't exist")
        getCaption(args.image, args.bearer, args.output)

    elif args.dir and args.bearer and args.output:
        if not os.path.exists(f"{args.dir}"):
            sys.exit(f"[!] {args.dir} doesn't exist")
        getCaptionsForDir(args.dir, args.bearer, args.output)

    elif args.pptx and args.update:
        if not os.path.exists(f"{args.update}") or not os.path.exists(f'{args.pptx}'):
            sys.exit(f"[!] {args.update} or {args.pptx} doesn't exist")
        updateImageAlttexts(args.pptx, args.update)

    elif args.pptx and args.list:
        if not os.path.exists(f'{args.pptx}'):
            sys.exit(f"[!] {args.pptx} doesn't exist")
        listImageAlttexts(args.pptx)

    else:
        print(f"[!] invalid...")
        parser.print_help()
